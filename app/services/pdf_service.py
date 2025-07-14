import os
import uuid
from PIL import Image  # Replaced fitz for image operations
import pdfplumber
from docx import Document
from pptx import Presentation
from pdf2docx import Converter as PdfToDocx
from pdf2image import convert_from_path
import zipfile
from io import BytesIO
from bs4 import BeautifulSoup
from tempfile import NamedTemporaryFile
from fastapi import UploadFile
import pdfkit
from pathlib import Path
import shutil
# from weasyprint import HTML
from PyPDF2 import PdfReader
from itertools import zip_longest
from pptx.util import Inches, Pt
import tempfile
import subprocess

class PDFConverter:
    
    async def _save_upload_file(self, file, extension):
        temp_file = NamedTemporaryFile(delete=False, suffix=f".{extension}")
        await file.seek(0)
        contents = await file.read()
        temp_file.write(contents)
        temp_file.close()
        return temp_file.name

    async def images_to_pdf(self, files: list[UploadFile]) -> str:
        """Convert multiple images to a single PDF using Pillow."""
        if not files:
            raise ValueError("No files provided for conversion.")

        output_path = f"temp_{uuid.uuid4()}.pdf"
        
        pil_images = []
        for file in files:
            image_bytes = await file.read()
            img = Image.open(BytesIO(image_bytes))
            
            # Convert to RGB to ensure compatibility for PDF saving.
            if img.mode == 'RGBA' or img.mode == 'P':
                img = img.convert('RGB')
                
            pil_images.append(img)
        
        if not pil_images:
             raise ValueError("No valid images could be processed.")

        # Save the first image, and append the rest
        pil_images[0].save(
            output_path, 
            "PDF", 
            resolution=100.0, 
            save_all=True, 
            append_images=pil_images[1:]
        )
        
        return output_path
    
    async def pdf_to_jpg(self, file: UploadFile, dpi: int = 300) -> str:
        """Convert each page of a PDF to a JPG image using pdf2image."""
        pdf_path = await self._save_upload_file(file, "pdf")
        
        try:
            # Convert PDF to a list of Pillow Image objects
            # This requires the poppler-utils package to be installed on the system
            images = convert_from_path(pdf_path, dpi=dpi, fmt='jpeg')
            
            if not images:
                raise RuntimeError("PDF to image conversion failed: no images were generated.")

            zip_buffer = BytesIO()
            with zipfile.ZipFile(zip_buffer, 'a', zipfile.ZIP_DEFLATED, False) as zip_file:
                for i, image in enumerate(images):
                    # Save image to a byte buffer
                    img_byte_arr = BytesIO()
                    image.save(img_byte_arr, format='JPEG')
                    img_byte_arr.seek(0)
                    
                    # Write image bytes to the zip file
                    zip_file.writestr(f'page_{i+1}.jpg', img_byte_arr.getvalue())
            
            output_path = f"temp_{uuid.uuid4()}.zip"
            with open(output_path, 'wb') as f:
                f.write(zip_buffer.getvalue())
            
            return output_path
        
        except Exception as e:
            # Provide a more helpful error if Poppler is likely missing
            if "pdftoppm" in str(e).lower() or "poppler" in str(e).lower():
                raise RuntimeError(
                    "pdf2image error: Poppler utility not found or failed. "
                    "Please ensure Poppler is installed in your Vercel environment."
                ) from e
            raise e
        finally:
            # Clean up the temporary PDF file
            os.unlink(pdf_path)

    
    async def word_to_pdf(self, file: UploadFile):
        """Convert Word to PDF with comprehensive error handling"""
        try:
            # 1. Create a safe working directory in system temp location
            temp_dir = os.path.join(tempfile.gettempdir(), f"conv_{uuid.uuid4().hex[:8]}")
            os.makedirs(temp_dir, exist_ok=True)
            
            # 2. Save input with simple filename
            input_path = os.path.join(temp_dir, "input.docx")
            with open(input_path, "wb") as f:
                await file.seek(0)
                f.write(await file.read())
            
            # 3. Prepare output path with simple name
            output_path = os.path.join(temp_dir, "output.pdf")
            
            # 4. Find LibreOffice
            soffice_path = self._find_soffice()
            if not soffice_path:
                return await self._fallback_word_to_pdf(file, temp_dir)

            # 5. Run conversion with absolute paths
            cmd = [
                soffice_path,
                '--headless',
                '--convert-to', 'pdf:writer_pdf_Export',
                '--outdir', temp_dir,
                os.path.abspath(input_path)
            ]
            
            result = subprocess.run(
                cmd,
                cwd=temp_dir,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=60
            )
            
            # 6. Check for output (LibreOffice sometimes changes the output filename)
            possible_outputs = [
                output_path,
                os.path.join(temp_dir, "input.pdf"),
                os.path.join(temp_dir, os.path.splitext(os.path.basename(input_path))[0] + ".pdf")
            ]
            
            final_output = None
            for possible_path in possible_outputs:
                if os.path.exists(possible_path):
                    final_output = possible_path
                    break
            
            if not final_output:
                # Try fallback if LibreOffice "succeeded" but didn't produce output
                return await self._fallback_word_to_pdf(file, temp_dir)
            
            # 7. Move result to final location
            result_path = os.path.join(tempfile.gettempdir(), f"result_{uuid.uuid4().hex[:8]}.pdf")
            shutil.move(final_output, result_path)
            
            return result_path
            
        except Exception as e:
            raise RuntimeError(f"Conversion failed: {str(e)}") from e
        finally:
            # Clean up
            if 'temp_dir' in locals() and os.path.exists(temp_dir):
                shutil.rmtree(temp_dir, ignore_errors=True)

    async def _fallback_word_to_pdf(self, file: UploadFile, temp_dir=None):
        """Fallback conversion using python-docx and reportlab"""
        try:
            if not temp_dir:
                temp_dir = os.path.join(tempfile.gettempdir(), f"fallback_{uuid.uuid4().hex[:8]}")
                os.makedirs(temp_dir, exist_ok=True)
            
            from docx import Document
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.units import inch
            
            # Save input
            docx_path = os.path.join(temp_dir, "input.docx")
            with open(docx_path, "wb") as f:
                await file.seek(0)
                f.write(await file.read())
            
            # Convert
            output_path = os.path.join(temp_dir, "output.pdf")
            doc = Document(docx_path)
            
            c = canvas.Canvas(output_path, pagesize=letter)
            y = 10 * inch  # Start position
            margin = 1 * inch
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:  # Skip empty paragraphs
                    if y < margin:  # New page if needed
                        c.showPage()
                        y = 10 * inch
                    c.setFont("Helvetica", 12)
                    c.drawString(margin, y, text)
                    y -= 0.25 * inch
            
            c.save()
            
            # Return final path
            result_path = os.path.join(tempfile.gettempdir(), f"result_{uuid.uuid4().hex[:8]}.pdf")
            shutil.move(output_path, result_path)
            return result_path
            
        except Exception as e:
            raise RuntimeError(f"Fallback conversion failed: {str(e)}") from e
        finally:
            if 'temp_dir' in locals() and os.path.exists(temp_dir):
                shutil.rmtree(temp_dir, ignore_errors=True)

    def _find_soffice(self):
        """Find LibreOffice with multiple fallback locations"""
        # Windows paths
        if os.name == 'nt':
            paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                "soffice.exe"
            ]
            # Check registry for install location
            try:
                import winreg
                with winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\LibreOffice\UNO\InstallPath") as key:
                    install_path = winreg.QueryValue(key, None)
                    paths.insert(0, os.path.join(install_path, "soffice.exe"))
            except:
                pass
        else:  # Linux/Mac
            paths = ["soffice", "/usr/bin/soffice", "/usr/local/bin/soffice"]
        
        for path in paths:
            try:
                if os.path.exists(path):
                    return path
                # Check PATH if not absolute path
                if not os.path.isabs(path):
                    found = shutil.which(path)
                    if found:
                        return found
            except:
                continue
    
        return None
    
    async def word_to_powerpoint(self, file: UploadFile):
        """Convert Word document to PowerPoint (text-only)"""
        docx_path = await self._save_upload_file(file, "docx")
        output_path = f"temp_{uuid.uuid4()}.pptx"
        
        try:
            from docx import Document
            from pptx import Presentation
            from pptx.util import Inches, Pt

            # Read the Word document
            doc = Document(docx_path)
            prs = Presentation()
            
            # Use a blank layout for all slides
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            
            # Process each paragraph in the Word doc
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:  # Skip empty paragraphs
                    # Create a new slide
                    slide = prs.slides.add_slide(blank_slide_layout)
                    
                    # Add a text box
                    left = Inches(1)
                    top = Inches(1)
                    width = prs.slide_width - Inches(2)
                    height = prs.slide_height - Inches(2)
                    
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = text_box.text_frame
                    text_frame.word_wrap = True
                    
                    # Add the paragraph text
                    p = text_frame.add_paragraph()
                    p.text = text
                    p.font.size = Pt(24)  # Standard font size for slides
            
            # Save the PowerPoint
            prs.save(output_path)
            return output_path
            
        except Exception as e:
            if os.path.exists(output_path):
                os.remove(output_path)
            raise RuntimeError(f"Word to PowerPoint conversion failed: {str(e)}")
        finally:
            os.unlink(docx_path)
            
    async def pdf_to_powerpoint(self, file: UploadFile) -> str:
        pdf_path = await self._save_upload_file(file, "pdf")
        output_path = f"temp_{uuid.uuid4()}.pptx"
        
        try:
            import aspose.pdf as ap
            from PyPDF2 import PdfReader, PdfWriter
            
            # Split the PDF into chunks of 4 pages each
            pdf_reader = PdfReader(pdf_path)
            total_pages = len(pdf_reader.pages)
            chunks = []
            
            # Create temporary directory for chunks
            temp_dir = os.path.join(tempfile.gettempdir(), f"pdf_chunks_{uuid.uuid4().hex[:8]}")
            os.makedirs(temp_dir, exist_ok=True)
            
            # Split PDF into 4-page chunks
            for i in range(0, total_pages, 4):
                chunk_path = os.path.join(temp_dir, f"chunk_{i//4}.pdf")
                pdf_writer = PdfWriter()
                
                end_page = min(i+4, total_pages)
                for page_num in range(i, end_page):
                    pdf_writer.add_page(pdf_reader.pages[page_num])
                
                with open(chunk_path, 'wb') as f:
                    pdf_writer.write(f)
                
                chunks.append(chunk_path)
            
            # Convert each chunk to PPTX using Aspose
            pptx_chunks = []
            for chunk in chunks:
                chunk_pptx = f"{chunk}.pptx"
                document = ap.Document(chunk)
                document.save(chunk_pptx, ap.PptxSaveOptions())
                pptx_chunks.append(chunk_pptx)
            
            # Combine all PPTX files into one
            if len(pptx_chunks) > 1:
                from pptx import Presentation
                
                # Create a new presentation
                combined_pptx = Presentation()
                
                # Add slides from each chunk
                for pptx_file in pptx_chunks:
                    source_pptx = Presentation(pptx_file)
                    for slide in source_pptx.slides:
                        # Create a new blank slide in the combined presentation
                        blank_slide_layout = combined_pptx.slide_layouts[6]  # Blank layout
                        new_slide = combined_pptx.slides.add_slide(blank_slide_layout)
                        
                        # Copy all shapes from source slide to new slide
                        for shape in slide.shapes:
                            new_shape = new_slide.shapes.add_shape(
                                shape.auto_shape_type,
                                shape.left, shape.top,
                                shape.width, shape.height
                            )
                            new_shape.text = shape.text
                            # Copy other properties as needed
                    
                    # Remove the temporary PPTX file
                    os.unlink(pptx_file)
                
                # Save the combined presentation
                combined_pptx.save(output_path)
                
                # Clean up remaining files
                for chunk in chunks:
                    os.unlink(chunk)
                os.rmdir(temp_dir)
            else:
                # If only one chunk, just rename it to output_path
                os.rename(pptx_chunks[0], output_path)
                os.rmdir(temp_dir)
            
            return output_path
        
        except Exception as e:
            # Clean up in case of error
            if 'temp_dir' in locals() and os.path.exists(temp_dir):
                for f in os.listdir(temp_dir):
                    os.unlink(os.path.join(temp_dir, f))
                os.rmdir(temp_dir)
            raise RuntimeError(f"PDF to PowerPoint conversion failed: {str(e)}")
        
        finally:
            if os.path.exists(pdf_path):
                os.unlink(pdf_path)
    
    async def html_to_pdf(self, html_content: str, is_url: bool = False):
        output_path = f"temp_{uuid.uuid4()}.pdf"
        try:
            # Try to find wkhtmltopdf in common locations
            wkhtmltopdf_path = None
            possible_paths = [
                '/usr/bin/wkhtmltopdf',       # Common Linux path
                '/usr/local/bin/wkhtmltopdf',  # Common macOS path
                shutil.which('wkhtmltopdf'),   # Check PATH environment variable
                'C:\\Program Files\\wkhtmltopdf\\bin\\wkhtmltopdf.exe'  # Windows default
            ]
            
            for path in possible_paths:
                if path and Path(path).exists():
                    wkhtmltopdf_path = path
                    break
            
            if not wkhtmltopdf_path:
                raise RuntimeError("wkhtmltopdf not found. Please install it first.")
            
            config = pdfkit.configuration(wkhtmltopdf=wkhtmltopdf_path)
            
            options = {
                'quiet': '',
                'page-size': 'A4',
                'encoding': "UTF-8",
            }

            if is_url:
                pdfkit.from_url(html_content, output_path, configuration=config, options=options)
            else:
                pdfkit.from_string(html_content, output_path, configuration=config, options=options)
            
            return output_path
        except Exception as e:
            if os.path.exists(output_path):
                os.remove(output_path)
            raise e
    async def pdf_to_word(self, file: UploadFile):
        pdf_path = await self._save_upload_file(file, "pdf")
        output_path = f"temp_{uuid.uuid4()}.docx"
        
        cv = PdfToDocx(pdf_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()
        
        os.unlink(pdf_path)
        return output_path

    async def powerpoint_to_pdf(self, file: UploadFile):
        """Convert PowerPoint to PDF with proper file handling"""
        temp_dir = None
        try:
            # Create temp directory
            temp_dir = os.path.join(tempfile.gettempdir(), f"pptx_pdf_{uuid.uuid4().hex[:8]}")
            os.makedirs(temp_dir, exist_ok=True)
            
            # Save input file
            pptx_path = os.path.join(temp_dir, "input.pptx")
            with open(pptx_path, "wb") as f:
                await file.seek(0)
                f.write(await file.read())
            
            # Use LibreOffice for conversion
            soffice_path = self._find_soffice()
            if not soffice_path:
                raise RuntimeError("LibreOffice not found. Please install it for conversion.")

            # Run conversion
            cmd = [
                soffice_path,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', temp_dir,
                pptx_path
            ]
            
            result = subprocess.run(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=60
            )
            
            # Find the output file
            output_pdf = None
            for fname in os.listdir(temp_dir):
                if fname.endswith('.pdf'):
                    output_pdf = os.path.join(temp_dir, fname)
                    break
            
            if not output_pdf:
                raise RuntimeError("Conversion failed - no PDF output created")
            
            # Create final filename (don't move, let FastAPI handle the file)
            final_filename = f"converted_{uuid.uuid4().hex[:8]}.pdf"
            final_path = os.path.join(tempfile.gettempdir(), final_filename)
            shutil.move(output_pdf, final_path)
            
            return final_path
            
        except Exception as e:
            if temp_dir and os.path.exists(temp_dir):
                shutil.rmtree(temp_dir, ignore_errors=True)
            raise RuntimeError(f"PPTX to PDF conversion failed: {str(e)}") from e
        
    async def pdf_to_excel(self ,file: UploadFile) -> str:
        pdf_path = await self._save_upload_file(file, "pdf")
        output_path = f"temp_{uuid.uuid4()}.xlsx"
        
        try:
            import aspose.pdf as ap
            document = ap.Document(pdf_path)
            save_options = ap.ExcelSaveOptions()
            document.save(output_path, save_options)
            return output_path
        finally:
            if os.path.exists(pdf_path):
                os.unlink(pdf_path)